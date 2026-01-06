#!/usr/bin/env python3
"""
EZ Platform Architecture Presentations Converter (Version 2)
Enhanced version with better PowerPoint support using SVG-to-PNG conversion.
"""

import os
import re
from pathlib import Path
from bs4 import BeautifulSoup
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
from PIL import Image
import cairosvg


class PresentationConverter:
    def __init__(self, base_dir):
        self.base_dir = Path(base_dir)
        self.output_dir = self.base_dir / "exports"
        self.output_dir.mkdir(exist_ok=True)

        # Create subdirectories for each format
        (self.output_dir / "pdf").mkdir(exist_ok=True)
        (self.output_dir / "markdown").mkdir(exist_ok=True)
        (self.output_dir / "powerpoint").mkdir(exist_ok=True)
        (self.output_dir / "images").mkdir(exist_ok=True)

    def get_html_files(self, lang_dir=None):
        """Get all HTML files to convert."""
        if lang_dir:
            search_dir = self.base_dir / lang_dir
        else:
            search_dir = self.base_dir

        html_files = []
        for file in search_dir.glob("*.html"):
            if file.name != "index.html":  # Skip index
                html_files.append(file)
        return html_files

    def extract_svg_from_html(self, html_path):
        """Extract SVG content from HTML file."""
        with open(html_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')

        # Extract title
        title = soup.find('h1')
        title_text = title.get_text() if title else html_path.stem

        # Extract subtitle
        subtitle = soup.find('p', class_='subtitle')
        subtitle_text = subtitle.get_text() if subtitle else ""

        # Extract SVG
        svg = soup.find('svg')
        svg_content = str(svg) if svg else None

        return {
            'title': title_text,
            'subtitle': subtitle_text,
            'svg': svg_content,
            'soup': soup
        }

    def svg_to_png(self, svg_content, output_path, width=1920, height=1080):
        """Convert SVG to PNG using cairosvg."""
        try:
            cairosvg.svg2png(
                bytestring=svg_content.encode('utf-8'),
                write_to=str(output_path),
                output_width=width,
                output_height=height,
                background_color='#1a1a2e'  # Match HTML background
            )
            return True
        except Exception as e:
            print(f"    WARNING: Could not convert SVG to PNG: {e}")
            return False

    def html_to_pdf(self, html_path, output_path):
        """Convert HTML to PDF using Playwright."""
        print(f"  Converting to PDF: {html_path.name} -> {output_path.name}")

        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page()

            # Load HTML file
            page.goto(f"file:///{html_path.absolute()}")

            # Wait for content to load
            page.wait_for_load_state('networkidle')

            # Generate PDF
            page.pdf(
                path=str(output_path),
                format='A4',
                landscape=True,
                print_background=True,
                margin={
                    'top': '0.5in',
                    'right': '0.5in',
                    'bottom': '0.5in',
                    'left': '0.5in'
                }
            )

            browser.close()

    def html_to_markdown_with_svg(self, html_path, output_path):
        """Convert HTML to Markdown with embedded SVG."""
        print(f"  Converting to Markdown: {html_path.name} -> {output_path.name}")

        content = self.extract_svg_from_html(html_path)

        markdown = []
        markdown.append(f"# {content['title']}\n")

        if content['subtitle']:
            markdown.append(f"_{content['subtitle']}_\n")

        markdown.append("\n---\n\n")

        if content['svg']:
            # Write SVG to separate file
            svg_filename = output_path.stem + ".svg"
            svg_path = output_path.parent / svg_filename

            with open(svg_path, 'w', encoding='utf-8') as f:
                f.write(content['svg'])

            # Embed SVG reference in markdown
            markdown.append(f"![{content['title']}](./{svg_filename})\n\n")

            # Also embed inline SVG for GitHub rendering
            markdown.append("## Architecture Diagram\n\n")
            markdown.append("```xml\n")
            markdown.append(content['svg'])
            markdown.append("\n```\n\n")

        # Extract any additional text content
        soup = content['soup']

        # Look for legend items
        legend = soup.find('div', class_='legend')
        if legend:
            markdown.append("## Legend\n\n")
            for item in legend.find_all('div', class_='legend-item'):
                text = item.get_text(strip=True)
                markdown.append(f"- {text}\n")
            markdown.append("\n")

        # Look for nav links
        nav = soup.find('div', class_='nav-links')
        if nav:
            markdown.append("## Related Documentation\n\n")
            for link in nav.find_all('a'):
                text = link.get_text(strip=True)
                href = link.get('href', '')
                if href and not href.startswith('http'):
                    # Convert HTML links to markdown links
                    md_href = href.replace('.html', '.md')
                    markdown.append(f"- [{text}](./{md_href})\n")
            markdown.append("\n")

        markdown.append("\n---\n\n")
        markdown.append(f"*Generated from {html_path.name}*\n")

        # Write markdown file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(''.join(markdown))

    def html_to_powerpoint(self, html_path, output_path):
        """Convert HTML to PowerPoint using SVG extraction."""
        print(f"  Converting to PowerPoint: {html_path.name} -> {output_path.name}")

        # Extract content
        content = self.extract_svg_from_html(html_path)

        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(13.33)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title.text = content['title']
        subtitle_shape.text = content['subtitle']

        # Diagram slide (if SVG exists)
        if content['svg']:
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)

            # Convert SVG to PNG
            png_path = self.output_dir / "images" / f"{output_path.stem}_diagram.png"
            if self.svg_to_png(content['svg'], png_path, width=1800, height=1200):
                # Add PNG to slide
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(12.33)

                try:
                    slide.shapes.add_picture(str(png_path), left, top, width=width)
                except Exception as e:
                    print(f"    WARNING: Could not add image to slide: {e}")

        # Save PowerPoint
        prs.save(str(output_path))

    def convert_file(self, html_path, lang_suffix=""):
        """Convert a single HTML file to all formats."""
        print(f"\nConverting: {html_path.name}")

        base_name = html_path.stem + lang_suffix

        # PDF
        pdf_output = self.output_dir / "pdf" / f"{base_name}.pdf"
        try:
            self.html_to_pdf(html_path, pdf_output)
        except Exception as e:
            print(f"    ERROR creating PDF: {e}")

        # Markdown with SVG
        md_output = self.output_dir / "markdown" / f"{base_name}.md"
        try:
            self.html_to_markdown_with_svg(html_path, md_output)
        except Exception as e:
            print(f"    ERROR creating Markdown: {e}")

        # PowerPoint
        pptx_output = self.output_dir / "powerpoint" / f"{base_name}.pptx"
        try:
            self.html_to_powerpoint(html_path, pptx_output)
        except Exception as e:
            print(f"    ERROR creating PowerPoint: {e}")

    def convert_all(self):
        """Convert all HTML files in the directory."""
        print("=" * 80)
        print("EZ Platform Architecture Presentations Converter (v2)")
        print("=" * 80)

        # Convert English files
        print("\n[English Presentations]")
        en_files = self.get_html_files()
        for html_file in en_files:
            self.convert_file(html_file)

        # Convert Hebrew files
        print("\n[Hebrew Presentations]")
        he_files = self.get_html_files("he")
        for html_file in he_files:
            self.convert_file(html_file, lang_suffix="_he")

        print("\n" + "=" * 80)
        print(f"Conversion complete!")
        print(f"Output directory: {self.output_dir}")
        print(f"  - PDF files: {self.output_dir / 'pdf'}")
        print(f"  - Markdown files: {self.output_dir / 'markdown'}")
        print(f"  - PowerPoint files: {self.output_dir / 'powerpoint'}")
        print(f"  - Images: {self.output_dir / 'images'}")
        print("=" * 80)


def main():
    # Get the directory containing this script
    script_dir = Path(__file__).parent

    # Create converter and run
    converter = PresentationConverter(script_dir)
    converter.convert_all()


if __name__ == "__main__":
    main()
